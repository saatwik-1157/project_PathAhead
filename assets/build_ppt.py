# -*- coding: utf-8 -*-
"""Build the PathAhead promotional deck (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image as PILImage
import os
IMG = os.path.join(os.path.dirname(os.path.abspath(__file__)), "img")

# ---- Brand palette ----
NAVY   = RGBColor(0x16, 0x29, 0x4A)
INK    = RGBColor(0x1E, 0x29, 0x3B)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
BLUE_D = RGBColor(0x1E, 0x40, 0xAF)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_D = RGBColor(0xEA, 0x58, 0x0C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF8, 0xFD)
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
SLATE  = RGBColor(0x1E, 0x29, 0x3B)
MUTED  = RGBColor(0x64, 0x74, 0x8B)
GREEN  = RGBColor(0x00, 0xD0, 0x84)
RED    = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x6A, 0x5C, 0xFF)
CYAN   = RGBColor(0x00, 0xC9, 0xFF)
# gradient stops
GRAD_DARK = [RGBColor(0x0D,0x16,0x33), RGBColor(0x1E,0x33,0x82), RGBColor(0x46,0x2C,0x9C)]  # navy > blue > purple
GRAD_BP   = [RGBColor(0x25,0x63,0xEB), RGBColor(0x6A,0x5C,0xFF)]   # blue > purple
GRAD_PC   = [RGBColor(0x6A,0x5C,0xFF), RGBColor(0x00,0xC9,0xFF)]   # purple > cyan
GRAD_BC   = [RGBColor(0x25,0x63,0xEB), RGBColor(0x00,0xC9,0xFF)]   # blue > cyan
GRAD_ORANGE = [RGBColor(0xF9,0x73,0x16), RGBColor(0xFB,0x92,0x3C)]
USE_GRAD = True
GRAD_MAP = { }  # filled after helpers

BODY = "Calibri"
HEAD = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    if bg == WHITE:
        set_gradient(r, [RGBColor(0xFF,0xFF,0xFF), RGBColor(0xEC,0xF1,0xFF)], 120)
    elif bg == LIGHT:
        set_gradient(r, [RGBColor(0xF4,0xF7,0xFE), RGBColor(0xE6,0xEE,0xFF)], 120)
    else:
        r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def set_gradient(shape, colors, angle_deg=120):
    spPr = shape._element.spPr
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn('a:gradFill'), {})
    gsLst = grad.makeelement(qn('a:gsLst'), {}); grad.append(gsLst)
    n = len(colors)
    for i, col in enumerate(colors):
        pos = int(round(i*100000/(n-1)))
        gs = grad.makeelement(qn('a:gs'), {'pos': str(pos)})
        srgb = grad.makeelement(qn('a:srgbClr'), {'val': str(col)})
        gs.append(srgb); gsLst.append(gs)
    lin = grad.makeelement(qn('a:lin'), {'ang': str(int(angle_deg*60000)), 'scaled': '1'})
    grad.append(lin)
    ln = spPr.find(qn('a:ln'))
    if ln is not None: ln.addprevious(grad)
    else: spPr.append(grad)
    try: shape.shadow.inherit = False
    except Exception: pass
    return shape

def slide_g(colors, angle=120):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_gradient(r, colors, angle)
    r.line.fill.background(); r.shadow.inherit = False
    return s

def soft_shadow(shape, blur=0.16, dist=0.07, alpha=17):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = spPr.makeelement(qn('a:effectLst'), {})
    shdw = eff.makeelement(qn('a:outerShdw'), {'blurRad': str(int(blur*914400)), 'dist': str(int(dist*914400)), 'dir': '5400000', 'rotWithShape': '0'})
    clr = shdw.makeelement(qn('a:srgbClr'), {'val': '0F172A'})
    a = clr.makeelement(qn('a:alpha'), {'val': str(int(alpha*1000))})
    clr.append(a); shdw.append(clr); eff.append(shdw)
    spPr.append(eff)
    return shape

def no_line(sh):
    sh.line.fill.background()

def shadow_off(sh):
    sh.shadow.inherit = False

def rect(s, x, y, w, h, fill=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None, line_w=1.0, radius=0.08):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    shadow_off(sp)
    # adjust corner radius for rounded rectangle
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    # premium: soft floating shadow on card-sized rounded rects
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE and w >= 2.0 and h >= 1.0 and fill is not None:
        soft_shadow(sp)
    return sp

def image_cover(s, name, x, y, w, h):
    """Place an image cropped to 'cover' the box (no distortion)."""
    path = os.path.join(IMG, name)
    iw, ih = PILImage.open(path).size
    img_ar = iw / ih; box_ar = w / h
    pic = s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if img_ar > box_ar:
        c = (1 - box_ar / img_ar) / 2.0
        pic.crop_left = c; pic.crop_right = c
    else:
        c = (1 - img_ar / box_ar) / 2.0
        pic.crop_top = c; pic.crop_bottom = c
    pic.shadow.inherit = False
    return pic

def overlay(s, x, y, w, h, color, alpha_pct):
    """Semi-transparent rectangle over an image (alpha_pct = opacity %)."""
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    srgb = sp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha_pct*1000))})
    srgb.append(a)
    return sp

GRAD_MAP = { str(BLUE): GRAD_BP, str(BLUE_D): GRAD_BC, str(PURPLE): GRAD_PC, str(NAVY): GRAD_BP }

def circle(s, x, y, d, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    if isinstance(fill, list):
        set_gradient(sp, fill, 120)
    elif USE_GRAD and str(fill) in GRAD_MAP:
        set_gradient(sp, GRAD_MAP[str(fill)], 120)
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    no_line(sp); shadow_off(sp)
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, bold, color, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, bold, color, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb

def P(txt, size, bold, color, font=BODY):
    return [(txt, size, bold, color, font)]

# ============================================================
# SLIDE 1 — COVER (dark)
# ============================================================
s = slide_g(GRAD_DARK, 120)
# students photo on the right
image_cover(s, "presentation.jpg", 8.55, 0, 4.783, 7.5)
# floating trust badge on the photo
rect(s, 8.95, 6.15, 3.95, 0.82, fill=WHITE, radius=0.24)
text(s, 8.95, 6.15, 3.95, 0.82, [ [("★ 4.9/5    ", 17, True, ORANGE, HEAD), ("Trusted by 30+ colleges", 12, True, NAVY, BODY)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# logo badge
badge = rect(s, 0.9, 0.8, 0.85, 0.85, fill=WHITE, radius=0.28)
text(s, 0.9, 0.8, 0.85, 0.85, [ [("PA", 24, True, NAVY, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.95, 0.9, 5, 0.7, [ [("PathAhead", 26, True, WHITE, HEAD)] ], anchor=MSO_ANCHOR.MIDDLE)
# eyebrow
rect(s, 0.9, 2.55, 4.35, 0.5, fill=RGBColor(0x22,0x38,0x5E), radius=0.5)
text(s, 0.9, 2.55, 4.35, 0.5, [ [("★  INDIA'S PREMIUM CAREER GUIDANCE", 11.5, True, RGBColor(0x9D,0xB8,0xEA), BODY)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# headline
text(s, 0.85, 3.25, 6.7, 2.0, [
    [("Your Career ", 46, True, WHITE, HEAD)],
    [("Starts Here.", 46, True, CYAN, HEAD)],
], space_after=0, line_spacing=1.0)
text(s, 0.9, 5.15, 6.5, 1.5, [
    [("Helping Intermediate students make smarter career decisions before college through practical workshops and future-ready learning.", 16, False, RGBColor(0xC7,0xD3,0xE8), BODY)],
], line_spacing=1.15)
text(s, 0.9, 6.85, 6.6, 0.5, [ [("Workshops for Inter 1st & 2nd Year Students", 12.5, True, RGBColor(0xF9,0x9A,0x4E), BODY)] ])

# ============================================================
# SLIDE 2 — THE PROBLEM (light)
# ============================================================
s = slide(WHITE)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("The problem we solve", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.7, [ [("Every year, thousands of students choose their future in the dark.", 18, False, MUTED, BODY)] ])
probs = [
    ("No real guidance", "Career and branch decisions made on hearsay, pressure or guesswork — not aptitude."),
    ("Unaware of new careers", "Little exposure to AI, Data Science and modern engineering domains before college."),
    ("Confusion over choices", "Overwhelmed by colleges, courses, scholarships and roadmaps with no clear map."),
    ("Decisions made too late", "By the time clarity comes, the best opportunities have already passed."),
]
x0, y0, cw, ch, gap = 0.9, 2.75, 5.7, 1.75, 0.35
for i,(t,d) in enumerate(probs):
    cx = x0 + (i%2)*(cw+gap); cy = y0 + (i//2)*(ch+gap)
    rect(s, cx, cy, cw, ch, fill=LIGHT, line=BORDER, line_w=1, radius=0.09)
    circle(s, cx+0.32, cy+0.32, 0.62, RED)
    text(s, cx+0.32, cy+0.32, 0.62, 0.62, [ [("×", 26, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+1.2, cy+0.28, cw-1.5, 0.5, [ [(t, 18, True, INK, HEAD)] ])
    text(s, cx+1.2, cy+0.78, cw-1.5, 0.9, [ [(d, 13.5, False, MUTED, BODY)] ], line_spacing=1.05)

# ============================================================
# SLIDE — FULL-BLEED PURPOSE QUOTE (cinematic)
# ============================================================
s = slide(NAVY)
image_cover(s, "students-group.jpg", 0, 0, 13.333, 7.5)
overlay(s, 0, 0, 13.333, 7.5, NAVY, 70)
text(s, 1.1, 2.2, 11.1, 0.5, [ [("OUR PURPOSE", 14, True, RGBColor(0xF9,0x9A,0x4E), BODY)] ])
text(s, 1.0, 2.8, 11.3, 2.2, [
    [("“No student should choose", 42, True, WHITE, HEAD)],
    [("their future by accident.”", 42, True, WHITE, HEAD)],
], line_spacing=1.05, space_after=0)
text(s, 1.1, 5.55, 10.8, 0.8, [ [("Give them clarity early — and you change the entire course of their life.", 18, False, RGBColor(0xD6,0xDF,0xEE), BODY)] ])

# ============================================================
# SLIDE 3 — MISSION & VISION (light)
# ============================================================
s = slide(LIGHT)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("Who we are", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("PathAhead brings honest, practical, future-ready career guidance directly into schools and junior colleges.", 17, False, MUTED, BODY)] ])
cards = [
    ("Our Mission", BLUE, "To empower every Intermediate student in India with the clarity, awareness and practical guidance they need to make confident, future-ready career decisions — before college."),
    ("Our Vision", ORANGE, "A future where career guidance is a standard part of every student's education — where choosing a path is driven by aptitude and awareness, not confusion or pressure."),
]
for i,(t,c,d) in enumerate(cards):
    cx = 0.9 + i*(5.85+0.4)
    rect(s, cx, 2.7, 5.85, 3.6, fill=CARD, line=BORDER, line_w=1, radius=0.06)
    circle(s, cx+0.5, 3.15, 0.75, c)
    text(s, cx+0.5, 3.15, 0.75, 0.75, [ [("✓", 30, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.5, 4.15, 4.9, 0.6, [ [(t, 24, True, INK, HEAD)] ])
    text(s, cx+0.5, 4.85, 4.9, 1.3, [ [(d, 15, False, MUTED, BODY)] ], line_spacing=1.12)

# ============================================================
# SLIDE 4 — OUR PROGRAMS (overview, light)
# ============================================================
s = slide(WHITE)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("Two signature workshops", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("Interactive 2-hour sessions, tailored to exactly where each student is in their journey.", 17, False, MUTED, BODY)] ])
progs = [
    ("FOR INTER 1st YEAR", BLUE, "Future Ready", "Lay a strong foundation early — career awareness, future skills and confidence from day one.",
     ["Career Awareness", "Engineering Branches", "AI Awareness", "Future Skills", "Communication Skills", "Personality Development", "Career Planning"]),
    ("FOR INTER 2nd YEAR", ORANGE, "Career Blueprint", "Make the right choices before college — a complete, decision-focused roadmap.",
     ["College Selection", "Technology Domains", "Programming Languages", "Scholarships", "Career Roadmap", "AI Careers", "Q&A Session"]),
]
for i,(tag,c,name,desc,items) in enumerate(progs):
    cx = 0.9 + i*(5.85+0.4)
    rect(s, cx, 2.55, 5.85, 4.5, fill=CARD, line=BORDER, line_w=1, radius=0.05)
    bar = rect(s, cx, 2.55, 5.85, 0.14, fill=c, shape=MSO_SHAPE.RECTANGLE)
    set_gradient(bar, GRAD_BP if c==BLUE else GRAD_ORANGE, 0)
    pill = rect(s, cx+0.5, 2.95, 3.2, 0.42, fill=(RGBColor(0xEF,0xF4,0xFF) if c==BLUE else RGBColor(0xFF,0xF3,0xEA)), radius=0.5)
    text(s, cx+0.5, 2.95, 3.2, 0.42, [ [(tag, 10.5, True, (BLUE if c==BLUE else ORANGE_D), BODY)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.5, 3.5, 4.9, 0.6, [ [(name, 27, True, INK, HEAD)] ])
    text(s, cx+0.5, 4.15, 4.9, 0.75, [ [(desc, 13.5, False, MUTED, BODY)] ], line_spacing=1.08)
    # two-column topic list
    col1 = items[:4]; col2 = items[4:]
    for j,it in enumerate(col1):
        text(s, cx+0.5, 5.0+j*0.42, 2.6, 0.4, [ [("✓  ", 12, True, c, BODY), (it, 12, False, SLATE, BODY)] ])
    for j,it in enumerate(col2):
        text(s, cx+3.15, 5.0+j*0.42, 2.5, 0.4, [ [("✓  ", 12, True, c, BODY), (it, 12, False, SLATE, BODY)] ])

# ============================================================
# SLIDE 5 — THE PATHAHEAD DIFFERENCE (comparison, light)
# ============================================================
s = slide(LIGHT)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("The PathAhead difference", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("We turn guesswork into a clear, confident plan.", 18, False, MUTED, BODY)] ])
# without
rect(s, 0.9, 2.7, 5.85, 3.9, fill=CARD, line=BORDER, line_w=1, radius=0.05)
circle(s, 1.35, 3.1, 0.55, RED); text(s, 1.35, 3.1, 0.55, 0.55, [ [("×", 22, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 2.1, 3.12, 4.4, 0.55, [ [("Without guidance", 20, True, INK, HEAD)] ], anchor=MSO_ANCHOR.MIDDLE)
without = ["Streams chosen on hearsay & pressure","Little awareness of AI & future careers","Confusion over colleges & scholarships","Decisions made too late, with regret"]
for j,it in enumerate(without):
    text(s, 1.35, 3.95+j*0.62, 5.1, 0.55, [ [("×   ", 15, True, RED, BODY), (it, 14.5, False, MUTED, BODY)] ], line_spacing=1.0)
# with
rect(s, 7.15, 2.7, 5.85, 3.9, fill=RGBColor(0xEF,0xF4,0xFF), line=BLUE, line_w=1.25, radius=0.05)
circle(s, 7.6, 3.1, 0.55, BLUE); text(s, 7.6, 3.1, 0.55, 0.55, [ [("✓", 22, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 8.35, 3.12, 4.4, 0.55, [ [("With PathAhead", 20, True, INK, HEAD)] ], anchor=MSO_ANCHOR.MIDDLE)
with_ = ["A clear, personalised sense of direction","Real awareness of AI & technology careers","Confidence in college & branch choices","A concrete roadmap they act on — now"]
for j,it in enumerate(with_):
    text(s, 7.6, 3.95+j*0.62, 5.1, 0.55, [ [("✓   ", 15, True, BLUE, BODY), (it, 14.5, False, SLATE, BODY)] ], line_spacing=1.0)

# ============================================================
# SLIDE 6 — CAREER DOMAINS (grid, light)
# ============================================================
s = slide(WHITE)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("Careers we open students' eyes to", 38, True, INK, HEAD)] ])
text(s, 0.9, 1.62, 11.4, 0.6, [ [("We make students aware of where tomorrow's opportunities really are.", 17, False, MUTED, BODY)] ])
domains = [
    ("AI", BLUE, "Artificial Intelligence", "Machine learning, generative AI and the careers reshaping every industry."),
    ("DS", ORANGE, "Data Science", "Turning data into decisions — one of today's most in-demand fields."),
    ("CS", PURPLE, "Cybersecurity", "Ethical hacking and security careers protecting the digital world."),
    ("CL", GREEN, "Cloud & DevOps", "The backbone of the internet and the fast-growing roles behind it."),
    ("SE", BLUE_D, "Software Engineering", "What developers actually do, and the languages worth learning first."),
    ("RB", ORANGE_D, "Robotics & IoT", "Where hardware meets intelligence — automation and mechatronics."),
]
x0,y0,cw,ch,gx,gy = 0.9, 2.5, 3.75, 1.95, 0.32, 0.3
for i,(ini,c,t,d) in enumerate(domains):
    cx = x0 + (i%3)*(cw+gx); cy = y0 + (i//3)*(ch+gy)
    rect(s, cx, cy, cw, ch, fill=LIGHT, line=BORDER, line_w=1, radius=0.09)
    circle(s, cx+0.3, cy+0.3, 0.7, c)
    text(s, cx+0.3, cy+0.3, 0.7, 0.7, [ [(ini, 15, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+1.15, cy+0.3, cw-1.35, 0.55, [ [(t, 15, True, INK, HEAD)] ])
    text(s, cx+0.3, cy+1.12, cw-0.6, 0.75, [ [(d, 11.5, False, MUTED, BODY)] ], line_spacing=1.05)

# ============================================================
# SLIDE 7 — WHAT STUDENTS RECEIVE (light)
# ============================================================
s = slide(LIGHT)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("What every student receives", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("A complete career starter kit they'll use all year long.", 17, False, MUTED, BODY)] ])
kit = [
    ("Career Guide","Clear guide to options after Intermediate, matched to interests."),
    ("Technology Roadmap","What to learn and when — a practical path through modern skills."),
    ("Engineering Guide","Every major branch explained — what it is and where it leads."),
    ("Career Checklist","A step-by-step checklist to keep career planning on track."),
    ("Resource Kit","Curated links, tools and reading to keep learning after."),
    ("Certificate + Community","Participation certificate and access to a future-ready community."),
]
x0,y0,cw,ch,gx,gy = 0.9, 2.55, 3.75, 1.9, 0.32, 0.3
cols = [BLUE, ORANGE, PURPLE, GREEN, BLUE_D, ORANGE_D]
for i,(t,d) in enumerate(kit):
    cx = x0 + (i%3)*(cw+gx); cy = y0 + (i//3)*(ch+gy)
    rect(s, cx, cy, cw, ch, fill=CARD, line=BORDER, line_w=1, radius=0.09)
    circle(s, cx+0.32, cy+0.32, 0.6, cols[i])
    text(s, cx+0.32, cy+0.32, 0.6, 0.6, [ [("✓", 20, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.3, cy+1.05, cw-0.55, 0.45, [ [(t, 15, True, INK, HEAD)] ])
    text(s, cx+0.3, cy+1.42, cw-0.55, 0.4, [ [(d, 11, False, MUTED, BODY)] ], line_spacing=1.02)

# ============================================================
# SLIDE 8 — FOR COLLEGES (why partner, light)
# ============================================================
s = slide(WHITE)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("Why colleges partner with us", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("A win for your students and your institution — with zero hassle for your team.", 17, False, MUTED, BODY)] ])
bens = [
    ("Better student outcomes","More informed stream, college and career decisions."),
    ("Enhanced reputation","Position your institution as forward-thinking and student-first."),
    ("Zero hassle","We bring content, mentors and materials — you provide the venue."),
    ("Future-ready content","AI, technology and modern engineering, kept current every term."),
    ("Flexible scheduling","Sessions arranged around your academic calendar."),
    ("Certificates & reports","Participation certificates and feedback to share with parents."),
]
x0,y0,cw,ch,gx,gy = 0.9, 2.55, 3.75, 1.85, 0.32, 0.3
for i,(t,d) in enumerate(bens):
    cx = x0 + (i%3)*(cw+gx); cy = y0 + (i//3)*(ch+gy)
    rect(s, cx, cy, cw, ch, fill=LIGHT, line=BORDER, line_w=1, radius=0.09)
    circle(s, cx+0.32, cy+0.3, 0.58, BLUE if i%2==0 else ORANGE)
    text(s, cx+0.32, cy+0.3, 0.58, 0.58, [ [("✓", 19, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.3, cy+1.0, cw-0.55, 0.45, [ [(t, 15, True, INK, HEAD)] ])
    text(s, cx+0.3, cy+1.37, cw-0.55, 0.4, [ [(d, 11, False, MUTED, BODY)] ], line_spacing=1.02)

# ============================================================
# SLIDE 9 — HOW IT WORKS (process, light)
# ============================================================
s = slide(LIGHT)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("How it works", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("From booking to breakthrough — a simple, structured, stress-free process.", 17, False, MUTED, BODY)] ])
steps = [
    ("1","Inquiry","You reach out via the form or phone with your requirements."),
    ("2","Planning","We confirm dates, batch size, stream and objectives together."),
    ("3","Delivery","A 2-hour interactive workshop delivered on your campus."),
    ("4","Follow-up","Certificates, resources and a feedback report for your records."),
]
x0,y0,cw = 0.9, 3.0, 2.86
gap = 0.28
for i,(n,t,d) in enumerate(steps):
    cx = x0 + i*(cw+gap)
    rect(s, cx, y0, cw, 2.7, fill=CARD, line=BORDER, line_w=1, radius=0.07)
    circle(s, cx+0.4, y0+0.4, 0.85, BLUE if i%2==0 else ORANGE)
    text(s, cx+0.4, y0+0.4, 0.85, 0.85, [ [(n, 30, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+0.4, y0+1.45, cw-0.7, 0.5, [ [(t, 19, True, INK, HEAD)] ])
    text(s, cx+0.4, y0+1.95, cw-0.7, 0.7, [ [(d, 12.5, False, MUTED, BODY)] ], line_spacing=1.05)

# ============================================================
# SLIDE — WORKSHOPS IN ACTION (photo gallery, light)
# ============================================================
s = slide(WHITE)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("PathAhead in action", 40, True, INK, HEAD)] ])
text(s, 0.9, 1.65, 11.4, 0.6, [ [("Interactive, engaging sessions delivered right on your campus.", 17, False, MUTED, BODY)] ])
gal = [("study.jpg","On-campus workshops"),("students-group.jpg","Engaged students"),("workshop.jpg","Career mentoring"),("tech.jpg","Future-ready skills")]
x0, y0 = 0.9, 2.6
gw, gh, gap = 2.83, 3.05, 0.29
for i,(img,cap) in enumerate(gal):
    cx = x0 + i*(gw+gap)
    image_cover(s, img, cx, y0, gw, gh)
    text(s, cx, y0+gh+0.12, gw, 0.4, [ [(cap, 12.5, True, INK, BODY)] ], align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10 — IMPACT / STATS (dark)
# ============================================================
s = slide(NAVY)
image_cover(s, "tech.jpg", 0, 0, 13.333, 7.5)
overlay(s, 0, 0, 13.333, 7.5, NAVY, 82)
text(s, 0.9, 0.85, 11.5, 0.9, [ [("Our impact so far", 40, True, WHITE, HEAD)] ])
text(s, 0.9, 1.8, 11, 0.6, [ [("Trusted across schools and junior colleges in India.", 17, False, RGBColor(0xC7,0xD3,0xE8), BODY)] ])
stats = [("5,000+","Students guided", ORANGE),("50+","Workshops delivered", CYAN),("30+","Partner institutions", ORANGE),("4.9/5","Average rating", CYAN)]
x0,y0,cw,gap = 0.9, 3.1, 2.86, 0.28
for i,(n,l,c) in enumerate(stats):
    cx = x0 + i*(cw+gap)
    rect(s, cx, y0, cw, 2.4, fill=RGBColor(0x1B,0x2E,0x52), radius=0.08)
    text(s, cx+0.2, y0+0.5, cw-0.4, 1.0, [ [(n, 48, True, c, HEAD)] ], align=PP_ALIGN.CENTER)
    text(s, cx+0.2, y0+1.6, cw-0.4, 0.6, [ [(l, 14, False, RGBColor(0xC7,0xD3,0xE8), BODY)] ], align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11 — TESTIMONIALS (light)
# ============================================================
s = slide(WHITE)
text(s, 0.9, 0.7, 11.5, 0.9, [ [("Trusted by students, teachers & principals", 34, True, INK, HEAD)] ])
quotes = [
    ("“The workshop gave me total clarity on which engineering branch fits me.”","Ananya R.","Inter 2nd Year Student", BLUE),
    ("“Engaging and relevant. Our students were talking about AI careers for days.”","S. Kumar","Faculty, Junior College", ORANGE),
    ("“Professional and valuable. We're making PathAhead an annual part of our calendar.”","Dr. Priya M.","Principal", NAVY),
]
x0,y0,cw,gap = 0.9, 2.35, 3.75, 0.32
for i,(q,n,role,c) in enumerate(quotes):
    cx = x0 + i*(cw+gap)
    rect(s, cx, y0, cw, 3.9, fill=LIGHT, line=BORDER, line_w=1, radius=0.06)
    text(s, cx+0.35, y0+0.2, 1, 0.8, [ [("“", 54, True, c, HEAD)] ])
    text(s, cx+0.35, y0+1.05, cw-0.7, 1.7, [ [(q, 15, False, SLATE, BODY)] ], line_spacing=1.15)
    circle(s, cx+0.35, y0+2.95, 0.6, c)
    text(s, cx+0.35, y0+2.95, 0.6, 0.6, [ [(n[0], 20, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, cx+1.1, y0+2.95, cw-1.3, 0.35, [ [(n, 14, True, INK, HEAD)] ])
    text(s, cx+1.1, y0+3.3, cw-1.3, 0.35, [ [(role, 11.5, False, MUTED, BODY)] ])

# ============================================================
# SLIDE 12 — CALL TO ACTION (dark)
# ============================================================
s = slide_g(GRAD_DARK, 120)
image_cover(s, "students-group.jpg", 8.55, 0, 4.783, 7.5)
text(s, 0.9, 1.05, 6.6, 1.5, [ [("Give your students an edge.", 38, True, WHITE, HEAD)] ], line_spacing=1.0)
text(s, 0.9, 2.55, 6.5, 1.0, [ [("Book a PathAhead workshop and help your students step into college with clarity and confidence.", 16, False, RGBColor(0xC7,0xD3,0xE8), BODY)] ], line_spacing=1.12)
contacts = [("Call us","+91 7780109877 • +91 9989057655"),("Email us","nirishapavuluri@gmail.com"),("","saathwik.13@gmail.com"),("Book online","Visit our website to book a workshop")]
y0=3.95
for i,(t,d) in enumerate(contacts):
    cy = y0 + i*0.72
    rect(s, 0.9, cy, 6.5, 0.6, fill=RGBColor(0x1B,0x2E,0x52), radius=0.16)
    circle(s, 1.1, cy+0.11, 0.38, BLUE if i%2==0 else ORANGE)
    text(s, 1.1, cy+0.11, 0.38, 0.38, [ [("✓", 13, True, WHITE, HEAD)] ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    lbl = (t+":  ") if t else ""
    text(s, 1.72, cy, 5.6, 0.6, [ [(lbl, 13, True, WHITE, BODY),(d, 12.5, False, RGBColor(0xC7,0xD3,0xE8), BODY)] ], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.9, 6.95, 6.6, 0.4, [ [("PathAhead — Your Career Starts Here.", 12.5, True, RGBColor(0xF9,0x9A,0x4E), BODY)] ])

out = r"C:\Users\Asus\project_nirdhan\assets\PathAhead-Presentation.pptx"
try:
    prs.save(out)
    print("Saved PathAhead-Presentation.pptx with", len(prs.slides._sldIdLst), "slides")
except PermissionError:
    alt = r"C:\Users\Asus\project_nirdhan\assets\PathAhead-Presentation-UPDATED.pptx"
    prs.save(alt)
    print("Main file was LOCKED (open in PowerPoint). Saved to PathAhead-Presentation-UPDATED.pptx instead.")
